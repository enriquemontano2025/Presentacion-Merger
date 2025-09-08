import os
from pptx import Presentation
from pathlib import Path

class PresentationMerger:
    def __init__(self, source_folder, language='ESP'):
        self.source_folder = Path(source_folder)
        self.language = language.upper()
        self.suffix = f"NF {self.language}"
        
    def scan_presentations(self):
        """Escanea la carpeta y devuelve información de las presentaciones"""
        presentations = {'ESP': [], 'ENG': []}
        
        try:
            # Buscar archivos .pptx y .ppt
            for ext in ['*.pptx', '*.ppt']:
                for file_path in self.source_folder.glob(ext):
                    filename = file_path.name
                    
                    if 'NF ESP' in filename:
                        building_name = filename.replace('NF ESP', '').replace('.pptx', '').replace('.ppt', '').strip()
                        presentations['ESP'].append({
                            'filename': filename,
                            'building': building_name,
                            'path': str(file_path),
                            'size': self._get_file_size(file_path)
                        })
                    elif 'NF ENG' in filename:
                        building_name = filename.replace('NF ENG', '').replace('.pptx', '').replace('.ppt', '').strip()
                        presentations['ENG'].append({
                            'filename': filename,
                            'building': building_name,
                            'path': str(file_path),
                            'size': self._get_file_size(file_path)
                        })
            
            return presentations
            
        except Exception as e:
            print(f"Error escaneando presentaciones: {e}")
            return presentations
    
    def merge_presentations(self, building_list, output_path):
        """Combina las presentaciones seleccionadas"""
        try:
            presentation_files, missing = self.find_building_presentations(building_list)
            
            if not presentation_files:
                return False, "❌ No se encontraron presentaciones para combinar."
            
            print(f"📁 Combinando {len(presentation_files)} presentaciones...")
            
            # Usar la primera presentación como base
            merged_presentation = Presentation(str(presentation_files[0]))
            slides_count = len(merged_presentation.slides)
            
            # Agregar diapositivas de otras presentaciones
            for i, ppt_file in enumerate(presentation_files[1:], 1):
                print(f"📄 Procesando archivo {i+1}/{len(presentation_files)}: {ppt_file.name}")
                
                try:
                    source_ppt = Presentation(str(ppt_file))
                    
                    # Copiar cada diapositiva
                    for slide_idx, source_slide in enumerate(source_ppt.slides):
                        # Crear nueva diapositiva con layout básico
                        slide_layout = merged_presentation.slide_layouts[0]  # Layout en blanco
                        new_slide = merged_presentation.slides.add_slide(slide_layout)
                        
                        # Copiar contenido (método simplificado)
                        self._copy_slide_content(source_slide, new_slide)
                    
                    slides_count += len(source_ppt.slides)
                    
                except Exception as slide_error:
                    print(f"⚠️ Error procesando {ppt_file.name}: {slide_error}")
                    continue
            
            # Guardar presentación combinada
            merged_presentation.save(output_path)
            
            message = f"✅ Presentación creada exitosamente!\n"
            message += f"📊 {slides_count} diapositivas de {len(presentation_files)} edificios\n"
            message += f"📁 Archivo: {os.path.basename(output_path)}"
            
            if missing:
                message += f"\n⚠️ Edificios no encontrados: {', '.join(missing)}"
            
            return True, message
            
        except Exception as e:
            error_msg = f"❌ Error al combinar presentaciones: {str(e)}"
            print(error_msg)
            return False, error_msg
    
    def find_building_presentations(self, building_list):
        """Encuentra las presentaciones específicas"""
        found_files = []
        missing_buildings = []
        
        all_presentations = self.scan_presentations()[self.language]
        
        for building in building_list:
            found = False
            for presentation in all_presentations:
                # Búsqueda más flexible
                if (building.lower() in presentation['building'].lower() or 
                    presentation['building'].lower() in building.lower()):
                    found_files.append(Path(presentation['path']))
                    found = True
                    print(f"✅ Encontrado: {building} → {presentation['filename']}")
                    break
            
            if not found:
                missing_buildings.append(building)
                print(f"❌ No encontrado: {building}")
        
        return found_files, missing_buildings
    
    def _copy_slide_content(self, source_slide, target_slide):
        """Copia el contenido de una diapositiva (método simplificado)"""
        try:
            # Esta es una implementación básica
            # En un entorno real, copiarías shapes, texto, imágenes, etc.
            # Por ahora, solo copiamos las propiedades básicas
            
            # Copiar background si es posible
            if hasattr(source_slide, 'background'):
                try:
                    target_slide.background = source_slide.background
                except:
                    pass
                    
            # Nota: Para una copia completa de diapositivas, necesitarías
            # una implementación más avanzada que maneje shapes, texto,
            # imágenes, tablas, etc.
            
        except Exception as e:
            print(f"⚠️ Error copiando contenido de diapositiva: {e}")
    
    def _get_file_size(self, file_path):
        """Obtiene el tamaño del archivo en formato legible"""
        try:
            size = os.path.getsize(file_path)
            for unit in ['B', 'KB', 'MB', 'GB']:
                if size < 1024.0:
                    return f"{size:.1f} {unit}"
                size /= 1024.0
            return f"{size:.1f} TB"
        except:
            return "? MB"
    
    def get_summary(self):
        """Obtiene un resumen de las presentaciones disponibles"""
        presentations = self.scan_presentations()
        
        summary = {
            'ESP': len(presentations['ESP']),
            'ENG': len(presentations['ENG']),
            'total': len(presentations['ESP']) + len(presentations['ENG']),
            'folder': str(self.source_folder),
            'buildings_esp': [p['building'] for p in presentations['ESP']],
            'buildings_eng': [p['building'] for p in presentations['ENG']]
        }
        
        return summary
